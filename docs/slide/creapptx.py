from pptx import Presentation
from pptx.util import Inches

def create_pptx():
    # Create a presentation object
    presentation = Presentation()

    # Slide 1: Title
    slide_layout = presentation.slide_layouts[0]  # Title slide layout
    title_slide = presentation.slides.add_slide(slide_layout)
    title = title_slide.shapes.title
    subtitle = title_slide.placeholders[1]

    title.text = "Server Code Overview"
    subtitle.text = "Socket-based multiplayer game server"

    # Slide 2: Code Overview
    slide_layout = presentation.slide_layouts[1]  # Title and content layout
    code_slide = presentation.slides.add_slide(slide_layout)
    title = code_slide.shapes.title
    content = code_slide.placeholders[1]

    title.text = "Code Overview"
    content.text = """
    import socket
    from Player import Player
    from Thread_game import game
    from Thread_waiting import waiting
    import threading
    import time
    import sqlite3

    # ... (rest of the code)
    """

    # Slide 3: Main Function
    slide_layout = presentation.slide_layouts[1]  # Title and content layout
    main_function_slide = presentation.slides.add_slide(slide_layout)
    title = main_function_slide.shapes.title
    content = main_function_slide.placeholders[1]

    title.text = "Main Function"
    content.text = """
    def main():
        # ... (rest of the main function)
    """

    # Slide 4: Server Configuration
    slide_layout = presentation.slide_layouts[1]  # Title and content layout
    config_slide = presentation.slides.add_slide(slide_layout)
    title = config_slide.shapes.title
    content = config_slide.placeholders[1]

    title.text = "Server Configuration"
    content.text = """
    # Configure the server
    global server_host
    global server_port
    server_host = '127.0.0.1'
    server_port = 12345
    """

    # ... (add more slides for other parts of the code)

    # Save the presentation to a file
    presentation.save("server_code_overview.pptx")

if __name__ == "__main__":
    create_pptx()
